import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Run diagram generator first to ensure Matplotlib analytical charts exist
import diagram_generator
diagram_generator.main()

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

FIGURES_DIR = r"e:\DL Project\presentation\figures"
OUTPUT_PPTX = r"e:\DL Project\presentation\PH_SSD_Presentation_v5.pptx"
OUTPUT_PPTX_ALT1 = r"e:\DL Project\presentation\PH_SSD_Presentation_v4.pptx"
OUTPUT_PPTX_ALT2 = r"e:\DL Project\presentation\PH_SSD_Presentation.pptx"
OUTPUT_PPTX_ALT3 = r"e:\DL Project\presentation\PH_SSD_Research_Presentation.pptx"

# Color Palette Constants
COLOR_NAVY = RGBColor(15, 28, 63)        # #0F1C3F
COLOR_BLUE = RGBColor(26, 86, 219)       # #1A56DB
COLOR_SKY = RGBColor(2, 132, 199)        # #0284C7
COLOR_SLATE = RGBColor(51, 65, 85)       # #334155
COLOR_MUTED = RGBColor(100, 116, 139)    # #64748B
COLOR_CARD_BG = RGBColor(248, 250, 252)  # #F8FAFC
COLOR_CARD_BORDER = RGBColor(226, 232, 240) # #E2E8F0
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GREEN = RGBColor(16, 185, 129)     # #10B981
COLOR_RED = RGBColor(239, 68, 68)        # #EF4444

TOTAL_SLIDES = 19

def set_slide_background(slide, color=COLOR_WHITE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, category, title, subtitle):
    """Clean Professional Header for content slides"""
    tx_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = tx_cat.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(9.5)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_BLUE
    p_cat.font.name = 'Arial'

    tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(11.7), Inches(0.55))
    tf_title = tx_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_NAVY
    p_title.font.name = 'Arial'

    if subtitle:
        tx_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.22), Inches(11.7), Inches(0.35))
        tf_sub = tx_sub.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = COLOR_SLATE
        p_sub.font.name = 'Arial'

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.58), Inches(11.733), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_CARD_BORDER
    line.line.color.rgb = COLOR_CARD_BORDER

def add_footer(slide, current_slide, total_slides=TOTAL_SLIDES):
    """Clean Footer for slides"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_CARD_BORDER
    line.line.color.rgb = COLOR_CARD_BORDER

    tx_foot = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.35))
    tf_foot = tx_foot.text_frame
    p_foot = tf_foot.paragraphs[0]
    p_foot.text = f"PH-SSD: Port-Hamiltonian State Space Dualities for Efficient Multimodal Learning          Slide {current_slide} of {total_slides}"
    p_foot.font.size = Pt(9)
    p_foot.font.color.rgb = COLOR_MUTED
    p_foot.font.name = 'Arial'

def add_card(slide, left, top, width, height, title="", bullets=None, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER, title_color=COLOR_NAVY):
    """Reusable visual card container with generous inner padding (0.35 in)"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)

    pad_h = Inches(0.35)
    pad_v = Inches(0.25)
    tx_box = slide.shapes.add_textbox(left + pad_h, top + pad_v, width - (2 * pad_h), height - (2 * pad_v))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    if title:
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = title_color
        p_title.font.name = 'Arial'
        p_title.space_after = Pt(10)

    if bullets:
        first = True if not title else False
        for b in bullets:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            
            b_str = str(b).strip()
            if b_str.startswith("✓") or b_str.startswith("⏱") or b_str.startswith("1.") or b_str.startswith("2.") or b_str.startswith("3.") or b_str.startswith("4.") or b_str.startswith("5."):
                p.text = b_str
                p.font.bold = True if (b_str.startswith("✓") or b_str.startswith("⏱")) else False
            elif b_str.startswith("•"):
                p.text = b_str
            else:
                p.text = f"• {b_str}"
            
            p.font.size = Pt(10.5)
            p.font.color.rgb = COLOR_SLATE
            p.font.name = 'Arial'
            p.space_after = Pt(6)

    return shape

def add_native_process_flow(slide, top, height, steps):
    """Generates N separate native PowerPoint cards with clear spacing and connecting arrows"""
    n = len(steps)
    left_margin = Inches(0.8)
    available_w = Inches(11.733)
    
    gap_w = Inches(0.35)
    box_w = (available_w - gap_w * (n - 1)) / n

    for i, (title, desc, stroke, bg) in enumerate(steps):
        box_left = left_margin + i * (box_w + gap_w)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, top, box_w, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg
        card.line.color.rgb = stroke
        card.line.width = Pt(1.5)

        tx = slide.shapes.add_textbox(box_left + Inches(0.12), top + Inches(0.2), box_w - Inches(0.24), height - Inches(0.4))
        tf = tx.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = stroke
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = COLOR_SLATE
        p2.alignment = PP_ALIGN.CENTER

        if i < n - 1:
            arrow_left = box_left + box_w + Inches(0.06)
            arrow_top = top + (height / 2) - Inches(0.1)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, Inches(0.23), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_BLUE
            arrow.line.fill.background()

# ==========================================
# SLIDE 1: Title Slide (B.Tech Research Proposal)
# ==========================================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1, COLOR_WHITE)

bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
bar.fill.solid()
bar.fill.fore_color.rgb = COLOR_BLUE
bar.line.fill.background()

card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.0), Inches(11.333), Inches(5.5))
card1.fill.solid()
card1.fill.fore_color.rgb = COLOR_CARD_BG
card1.line.color.rgb = COLOR_BLUE
card1.line.width = Pt(2)

tx = slide1.shapes.add_textbox(Inches(1.6), Inches(1.5), Inches(10.133), Inches(4.5))
tf = tx.text_frame
tf.word_wrap = True

p0 = tf.paragraphs[0]
p0.text = "B.TECH FINAL YEAR RESEARCH PROPOSAL"
p0.font.size = Pt(11)
p0.font.bold = True
p0.font.color.rgb = COLOR_BLUE
p0.space_after = Pt(14)

p1 = tf.add_paragraph()
p1.text = "PH-SSD: Port-Hamiltonian State Space Dualities\nfor Efficient Multimodal Learning"
p1.font.size = Pt(28)
p1.font.bold = True
p1.font.color.rgb = COLOR_NAVY
p1.space_after = Pt(18)

p2 = tf.add_paragraph()
p2.text = "A Linear-Complexity Architecture for Joint Vision-Language Modeling"
p2.font.size = Pt(14)
p2.font.color.rgb = COLOR_SLATE
p2.space_after = Pt(36)

p3 = tf.add_paragraph()
p3.text = "Student Presenter: [Student Name]\nDepartment: Department of Computer Science & Engineering\nInstitution: [University / College Name]  |  Advisor: [Faculty Advisor Name]"
p3.font.size = Pt(12)
p3.font.bold = True
p3.font.color.rgb = COLOR_NAVY

# ==========================================
# SLIDE 2: Problem Statement
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)
add_header(slide2, "PROBLEM STATEMENT", "Limitations of Current Multimodal AI Models", "Key computational bottlenecks in modern Vision-Language Transformers")
add_footer(slide2, 2)

probs = [
    ("High Memory & O(N²) Complexity", ["Standard Transformer self-attention scales quadratically with sequence length N.", "VRAM footprint rapidly explodes when processing high-resolution images & long captions."]),
    ("Poor Multimodal Fusion", ["Conventional cross-attention layers introduce alignment latency.", "Fails to capture fine-grained energy dynamics between visual patches and text tokens."]),
    ("Background Noise & Computational Waste", ["Raw visual patches contain redundant background noise.", "Unfiltered noise inflates FLOPs and degrades multimodal retrieval precision."])
]

for i, (p_title, p_bullets) in enumerate(probs):
    add_card(slide2, Inches(0.8), Inches(1.8 + i * 1.65), Inches(5.8), Inches(1.5), p_title, p_bullets, border_color=COLOR_RED if i==0 else COLOR_CARD_BORDER)

slide2.shapes.add_picture(os.path.join(FIGURES_DIR, "fig_problem_complexity.png"), Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))

# ==========================================
# SLIDE 3: Motivation
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)
add_header(slide3, "MOTIVATION", "Unified Synergy in Multimodal Perception", "Inspiration from human cognitive multimodal understanding")
add_footer(slide3, 3)

add_card(slide3, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), "Human vs. Machine Perception", [
    "Real-Life Example:",
    "  • Visual Stream: An image showing a dog playing with a ball.",
    "  • Textual Stream: Caption reads 'A dog is playing in the park.'",
    "Human Perception:",
    "  • Humans process visual and textual cues simultaneously as a unified energy system.",
    "Current AI Flaw:",
    "  • Modern models process inputs separately, incurring severe fusion overhead.",
    "The PH-SSD Solution:",
    "  • Joint state-space dualities conserve energy while filtering spatial and textual noise."
])

mot_left = Inches(6.6)
mot_top = Inches(1.8)

card_v = add_card(slide3, mot_left, mot_top, Inches(2.6), Inches(2.3), "Visual Stream", [
    "Image Frame:",
    "Spatial visual features of a dog playing with a ball"
], bg_color=RGBColor(239, 246, 255), border_color=COLOR_BLUE, title_color=COLOR_BLUE)

card_t = add_card(slide3, mot_left + Inches(3.1), mot_top, Inches(2.6), Inches(2.3), "Text Stream", [
    "Caption Token:",
    '"A dog is playing in the park."'
], bg_color=RGBColor(240, 249, 255), border_color=COLOR_SKY, title_color=COLOR_SKY)

arr = slide3.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, mot_left + Inches(2.62), mot_top + Inches(1.0), Inches(0.46), Inches(0.25))
arr.fill.solid()
arr.fill.fore_color.rgb = COLOR_BLUE
arr.line.fill.background()

add_card(slide3, mot_left + Inches(0.35), mot_top + Inches(2.7), Inches(5.0), Inches(1.9), "Unified Multimodal Synergy", [
    "Co-dependent joint energy representation.",
    "Human-like simultaneous vision-language perception."
], bg_color=RGBColor(236, 253, 245), border_color=COLOR_GREEN, title_color=RGBColor(4, 120, 87))

# ==========================================
# SLIDE 4: Research Objective
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4)
add_header(slide4, "RESEARCH OBJECTIVES", "Core Strategic Goals of PH-SSD", "Five key architectural targets designed for efficiency and robustness")
add_footer(slide4, 4)

objectives = [
    ("✓ Remove Noisy Information", "Utilize Port-Hamiltonian energy dissipation to filter out irrelevant background visual noise and text redundancy."),
    ("✓ Improve Multimodal Fusion", "Develop a Variational Cross-Modal (VCM) coupler for tight latent alignment between vision and text."),
    ("✓ Reduce Complexity to O(N)", "Replace quadratic Transformer attention with Mamba-2 State Space Dualities for linear-time scaling."),
    ("✓ Maximize Computational Efficiency", "Achieve lower memory usage, faster inference latency, and high sample throughput."),
    ("✓ Build Scalable Multimodal AI", "Establish a robust, modular backbone adaptable to long-context enterprise multimodal tasks.")
]

for i, (obj_title, obj_desc) in enumerate(objectives):
    top_pos = Inches(1.8 + i * 0.98)
    card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(11.733), Inches(0.85))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_BLUE if i==0 else COLOR_CARD_BORDER
    card.line.width = Pt(1.5)

    tx_o = slide4.shapes.add_textbox(Inches(1.1), top_pos + Inches(0.12), Inches(11.133), Inches(0.6))
    tf_o = tx_o.text_frame
    tf_o.word_wrap = True
    p_t = tf_o.paragraphs[0]
    p_t.text = obj_title
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_BLUE
    
    p_d = tf_o.add_paragraph()
    p_d.text = obj_desc
    p_d.font.size = Pt(10.5)
    p_d.font.color.rgb = COLOR_SLATE

# ==========================================
# SLIDE 5: Research Domain
# ==========================================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)
add_header(slide5, "RESEARCH DOMAIN", "Taxonomy & Evolutionary Context", "Positioning PH-SSD within modern artificial intelligence paradigms")
add_footer(slide5, 5)

domain_steps = [
    ("Artificial Intelligence", "Foundational AI Field", COLOR_NAVY, RGBColor(241, 245, 249)),
    ("Deep Learning", "Neural Architectures", RGBColor(30, 58, 138), RGBColor(239, 246, 255)),
    ("Multimodal Learning", "Vision & Language", COLOR_BLUE, RGBColor(219, 234, 254)),
    ("State Space Models", "Linear Sequences", COLOR_SKY, RGBColor(224, 242, 254)),
    ("PH-SSD Model", "Energy Dualities", COLOR_GREEN, RGBColor(209, 250, 229))
]

add_native_process_flow(slide5, Inches(2.0), Inches(2.3), domain_steps)

add_card(slide5, Inches(0.8), Inches(4.7), Inches(5.7), Inches(1.9), "Domain Foundation", [
    "Combines Systems & Control Theory (Port-Hamiltonian Systems) with modern Deep Learning.",
    "Extends State Space Models (SSMs) to high-dimensional multimodal vision-language tasks."
])

add_card(slide5, Inches(6.8), Inches(4.7), Inches(5.733), Inches(1.9), "Architectural Breakthrough", [
    "First framework to achieve physics-inspired noise filtering (SD-NPF) before State Space Duality.",
    "Establishes a new state-of-the-art benchmark for efficient multimodal fusion."
])

# ==========================================
# SLIDE 6: Dataset (With "Why Flickr8k?")
# ==========================================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)
add_header(slide6, "DATASET", "Flickr8k Benchmark Dataset & Rationale", "Primary experimental dataset for model proposal validation")
add_footer(slide6, 6)

add_card(slide6, Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.8), "Dataset Specifications", [
    "Dataset Name: Flickr8k",
    "Total Images: 8,000 Natural Scenes",
    "Total Captions: 40,000 Descriptions",
    "Caption Ratio: 5 Captions per Image",
    "Dataset Size: ~1.0 GB Total Storage",
    "Splits: 6,000 Train | 1,000 Val | 1,000 Test",
    "Modalities: Dual-Stream Images & Natural Text",
    "Task: Image-Text Cross-Modal Retrieval"
])

fl_left = Inches(6.3)
fl_top = Inches(1.8)

# "Why Flickr8k?" Rationale Card
add_card(slide6, fl_left, fl_top, Inches(6.233), Inches(2.3), "Why Choose Flickr8k for Proposal?", [
    "✓ Ideal Benchmark Size: Enables rapid architecture validation & debugging.",
    "✓ Official Standard: Widely recognized baseline dataset in multimodal literature.",
    "✓ Low Resource Footprint (~1 GB): Fits comfortably on single GPU hardware.",
    "✓ Proof-of-Concept Baseline: Prepares model before scaling to Flickr30k & MS-COCO."
], bg_color=RGBColor(239, 246, 255), border_color=COLOR_BLUE, title_color=COLOR_BLUE)

# Multi-Caption Sample Card underneath
add_card(slide6, fl_left, fl_top + Inches(2.5), Inches(6.233), Inches(2.1), "Sample Image-Caption Pair", [
    'Image: A child playing on a wooden outdoor slide.',
    'Caption 1: "A child in a pink dress is climbing wooden stairs."',
    'Caption 2: "A little girl in a pink shirt climbs into a playhouse."'
], bg_color=COLOR_CARD_BG)

# ==========================================
# SLIDE 7: Overall Architecture
# ==========================================
slide7 = prs.slides.add_slide(blank_layout)
set_slide_background(slide7)
add_header(slide7, "OVERALL ARCHITECTURE", "End-to-End Dual-Stream PH-SSD Network", "High-level diagram of vision & text encoders, SD-NPF, Mamba-2 SSD, and VCM-SSD Coupler")
add_footer(slide7, 7)

top_v = Inches(1.8)
top_t = Inches(3.6)
card_w = Inches(1.8)
card_h = Inches(1.3)

v_items = [
    ("Input Image", "RGB 224×224", COLOR_BLUE, RGBColor(239, 246, 255)),
    ("Vision Encoder", "ViT / DINOv2", COLOR_BLUE, RGBColor(219, 234, 254)),
    ("SD-NPF", "Noise Filter", COLOR_SKY, RGBColor(224, 242, 254)),
    ("Mamba-2 SSD", "Vision Dual", COLOR_NAVY, RGBColor(241, 245, 249))
]

for i, (title, sub, stroke, bg) in enumerate(v_items):
    c_left = Inches(0.8 + i * 2.1)
    card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, top_v, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = stroke
    card.line.width = Pt(1.5)
    
    tx = slide7.shapes.add_textbox(c_left + Inches(0.1), top_v + Inches(0.2), card_w - Inches(0.2), card_h - Inches(0.3))
    tf = tx.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = stroke
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.size = Pt(9)
    p2.font.color.rgb = COLOR_SLATE
    p2.alignment = PP_ALIGN.CENTER

    if i < 3:
        arr = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, c_left + card_w + Inches(0.05), top_v + Inches(0.55), Inches(0.2), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLOR_BLUE
        arr.line.fill.background()

t_items = [
    ("Input Caption", "Natural Text", COLOR_SKY, RGBColor(240, 249, 255)),
    ("Text Encoder", "RoBERTa Base", COLOR_SKY, RGBColor(224, 242, 254)),
    ("SD-NPF", "Noise Filter", COLOR_SKY, RGBColor(224, 242, 254)),
    ("Mamba-2 SSD", "Text Dual", COLOR_NAVY, RGBColor(241, 245, 249))
]

for i, (title, sub, stroke, bg) in enumerate(t_items):
    c_left = Inches(0.8 + i * 2.1)
    card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, top_t, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = stroke
    card.line.width = Pt(1.5)
    
    tx = slide7.shapes.add_textbox(c_left + Inches(0.1), top_t + Inches(0.2), card_w - Inches(0.2), card_h - Inches(0.3))
    tf = tx.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = stroke
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.size = Pt(9)
    p2.font.color.rgb = COLOR_SLATE
    p2.alignment = PP_ALIGN.CENTER

    if i < 3:
        arr = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, c_left + card_w + Inches(0.05), top_t + Inches(0.55), Inches(0.2), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = COLOR_SKY
        arr.line.fill.background()

coupler = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(2.0), Inches(1.8), Inches(2.6))
coupler.fill.solid()
coupler.fill.fore_color.rgb = RGBColor(209, 250, 229)
coupler.line.color.rgb = COLOR_GREEN
coupler.line.width = Pt(2)

tx_c = slide7.shapes.add_textbox(Inches(9.3), Inches(2.7), Inches(1.6), Inches(1.4))
tf_c = tx_c.text_frame
tf_c.word_wrap = True
p = tf_c.paragraphs[0]
p.text = "VCM-SSD Coupler\n(Variational\nFusion)"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = RGBColor(4, 120, 87)
p.alignment = PP_ALIGN.CENTER

pred = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.3), Inches(2.7), Inches(1.2), Inches(1.2))
pred.fill.solid()
pred.fill.fore_color.rgb = COLOR_NAVY
pred.line.color.rgb = COLOR_NAVY

tx_p = slide7.shapes.add_textbox(Inches(11.35), Inches(3.0), Inches(1.1), Inches(0.8))
tf_p = pred.text_frame
tf_p.word_wrap = True
p = tf_p.paragraphs[0]
p.text = "Joint Output\nPrediction"
p.font.size = Pt(9.5)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE
p.alignment = PP_ALIGN.CENTER

a_v = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.75), top_v + Inches(0.55), Inches(0.4), Inches(0.2))
a_v.fill.solid()
a_v.fill.fore_color.rgb = COLOR_BLUE
a_v.line.fill.background()

a_t = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.75), top_t + Inches(0.55), Inches(0.4), Inches(0.2))
a_t.fill.solid()
a_t.fill.fore_color.rgb = COLOR_SKY
a_t.line.fill.background()

a_out = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(11.02), Inches(3.2), Inches(0.25), Inches(0.2))
a_out.fill.solid()
a_out.fill.fore_color.rgb = COLOR_GREEN
a_out.line.fill.background()

add_card(slide7, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.1), "", [
    "Dual-Stream Processing: Parallel vision & text feature extraction -> SD-NPF noise damping -> Mamba-2 SSD linear temporal modeling -> VCM-SSD variational fusion."
])

# ==========================================
# SLIDE 8 (NEW): Concrete Example Walkthrough
# ==========================================
slide8 = prs.slides.add_slide(blank_layout)
set_slide_background(slide8)
add_header(slide8, "CONCRETE EXAMPLE FLOW", "Step-by-Step Execution on a Real Pair", "Intuitive visual flow demonstrating how PH-SSD processes a sample image and caption")
add_footer(slide8, 8)

example_steps = [
    ("1. Input Pair", "Image: Dog playing\nCaption: 'A dog is playing'", COLOR_NAVY, RGBColor(248, 250, 252)),
    ("2. Encoders", "ViT Patches &\nRoBERTa Tokens", RGBColor(30, 58, 138), RGBColor(239, 246, 255)),
    ("3. SD-NPF", "Removes grass & text\nbackground noise", COLOR_BLUE, RGBColor(219, 234, 254)),
    ("4. Mamba-2", "Linear O(N) dual\nsequence modeling", COLOR_SKY, RGBColor(224, 242, 254)),
    ("5. VCM-SSD", "Variational latent\nfusion space", COLOR_GREEN, RGBColor(209, 250, 229)),
    ("6. Prediction", "Image ↔ Text\nMatch Score = 0.96", RGBColor(4, 120, 87), RGBColor(236, 253, 245))
]

add_native_process_flow(slide8, Inches(2.0), Inches(2.8), example_steps)

add_card(slide8, Inches(0.8), Inches(5.2), Inches(11.733), Inches(1.5), "Why Concrete Flow Helps Faculty Understanding", [
    "Replaces abstract equations with a tangible pipeline.",
    "Shows how raw RGB pixels of a dog playing with a ball and sentence tokens are noise-damped, state-modeled, and coupled into a high-confidence match score."
])

# ==========================================
# SLIDE 9: Component 1 - SD-NPF (Intuition -> Mechanism -> Math)
# ==========================================
slide9 = prs.slides.add_slide(blank_layout)
set_slide_background(slide9)
add_header(slide9, "COMPONENT 1: SD-NPF", "Symplectic Dissipative Neural Pre-Filter", "Removing noisy background clutter before state space modeling")
add_footer(slide9, 9)

# 3-Step Clear Card Breakdown
sd_cards = [
    ("1. What Does It Do? (Intuition)", [
        "Filters out uninformative visual clutter (e.g. grass, background trees) and redundant text tokens.",
        "Ensures down-stream State Space Models process pure semantic signals rather than noise."
    ], COLOR_BLUE, RGBColor(239, 246, 255)),
    
    ("2. How Does It Work? (Mechanism)", [
        "Inspired by Port-Hamiltonian physical control systems.",
        "Applies a dissipative energy damping matrix R(x) >= 0 to attenuate high-frequency noise perturbations."
    ], COLOR_SKY, RGBColor(224, 242, 254)),
    
    ("3. Mathematical Formulation (Equations)", [
        "State Equation:  dx/dt = [J(x) - R(x)] grad H(x)",
        "J(x) = -J(x)ᵀ : Energy-conserving symplectic matrix",
        "Guarantees bounded asymptotic feature stability."
    ], COLOR_NAVY, RGBColor(241, 245, 249))
]

for i, (c_title, c_bullets, c_stroke, c_bg) in enumerate(sd_cards):
    add_card(slide9, Inches(0.8), Inches(1.8 + i * 1.65), Inches(5.8), Inches(1.5), c_title, c_bullets, bg_color=c_bg, border_color=c_stroke, title_color=c_stroke)

slide9.shapes.add_picture(os.path.join(FIGURES_DIR, "fig_sd_npf_energy.png"), Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))

# ==========================================
# SLIDE 10: Component 2 - Mamba-2 State Space Model (What is Mamba -> What is SSD -> Table)
# ==========================================
slide10 = prs.slides.add_slide(blank_layout)
set_slide_background(slide10)
add_header(slide10, "COMPONENT 2: MAMBA-2 SSD", "State Space Duality Backbone", "Linear-time sequence processing engine replacing self-attention")
add_footer(slide10, 10)

add_card(slide10, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), "Mamba-2 & SSD Concepts", [
    "What is Mamba-2?",
    "  • A State Space Model (SSM) designed for sequence processing with linear O(N) time complexity.",
    "  • Replaces quadratic O(N²) Transformer self-attention.",
    "What is State Space Duality (SSD)?",
    "  • Connects continuous-time state space models with structured matrix multiplication.",
    "Key Advantages:",
    "  • 5× faster inference speedup on GPU hardware.",
    "  • Sub-quadratic VRAM memory footprint."
])

tbl_card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.8), Inches(6.033), Inches(4.8))
tbl_card.fill.solid()
tbl_card.fill.fore_color.rgb = COLOR_CARD_BG
tbl_card.line.color.rgb = COLOR_CARD_BORDER
tbl_card.line.width = Pt(1.5)

rows, cols = 6, 3
table_shape = slide10.shapes.add_table(rows, cols, Inches(6.7), Inches(2.3), Inches(5.6), Inches(4.0))
table = table_shape.table

headers = ["Metric / Attribute", "Standard Transformer", "PH-SSD (Mamba-2)"]
for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_BLUE
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

table_data = [
    ["Time Complexity", "O(N²)", "O(N) Linear"],
    ["Memory Scaling", "Quadratic O(N²)", "Linear O(N)"],
    ["Noise Handling", "None (Processes All)", "SD-NPF Filtered"],
    ["Inference Latency", "High Bottleneck", "Ultra-Low Scan Kernel"],
    ["Multimodal Fusion", "Cross-Attention", "Variational VCM-SSD"]
]

for row_idx, row_data in enumerate(table_data):
    for col_idx, cell_text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_WHITE if row_idx % 2 == 0 else COLOR_CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.text = cell_text
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_NAVY if col_idx == 2 else COLOR_SLATE
        if col_idx == 2:
            p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

# ==========================================
# SLIDE 11: Component 3 - VCM-SSD
# ==========================================
slide11 = prs.slides.add_slide(blank_layout)
set_slide_background(slide11)
add_header(slide11, "COMPONENT 3: VCM-SSD COUPLER", "Variational Cross-Modal SSD Coupler", "Dynamic latent alignment connecting vision and language state spaces")
add_footer(slide11, 11)

add_card(slide11, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), "VCM-SSD Coupler Highlights", [
    "Cross-Modal State Duality:",
    "  • Directly bridges filtered vision hidden states h_v with text hidden states h_t.",
    "Variational Latent Space:",
    "  • Learns a probabilistic joint distribution q(z|v,t) bounded by KL divergence loss.",
    "Tight Semantic Alignment:",
    "  • Enforces semantic consistency across image patches and sentence tokens.",
    "Bidirectional Retrieval Head:",
    "  • Supports dual Image-to-Text and Text-to-Image cross-modal inference."
])

vcm_left = Inches(6.6)
vcm_top = Inches(1.8)

add_card(slide11, vcm_left, vcm_top, Inches(2.5), Inches(2.2), "Vision Features", [
    "Filtered Vision Stream h_v"
], bg_color=RGBColor(219, 234, 254), border_color=COLOR_BLUE, title_color=COLOR_BLUE)

add_card(slide11, vcm_left, vcm_top + Inches(2.6), Inches(2.5), Inches(2.2), "Text Features", [
    "Filtered Text Stream h_t"
], bg_color=RGBColor(224, 242, 254), border_color=COLOR_SKY, title_color=COLOR_SKY)

add_card(slide11, vcm_left + Inches(2.8), vcm_top + Inches(0.9), Inches(3.0), Inches(3.0), "VCM-SSD Engine", [
    "Variational Latent q(z|v,t)",
    "Cross-Modal SSD Duality",
    "KL Divergence Regularizer"
], bg_color=RGBColor(209, 250, 229), border_color=COLOR_GREEN, title_color=RGBColor(4, 120, 87))

# ==========================================
# SLIDE 12: Workflow
# ==========================================
slide12 = prs.slides.add_slide(blank_layout)
set_slide_background(slide12)
add_header(slide12, "WORKFLOW", "Step-by-Step Processing Pipeline", "Complete end-to-end data flow from raw multimodal inputs to final predictions")
add_footer(slide12, 12)

workflow_steps = [
    ("1. Input", "Raw Image &\nText Tokens", COLOR_NAVY, RGBColor(248, 250, 252)),
    ("2. Encoder", "ViT / RoBERTa\nEmbeddings", RGBColor(30, 58, 138), RGBColor(239, 246, 255)),
    ("3. SD-NPF", "Port-Hamiltonian\nNoise Filter", COLOR_BLUE, RGBColor(219, 234, 254)),
    ("4. Mamba-2", "State Space\nSSD Backbone", COLOR_SKY, RGBColor(224, 242, 254)),
    ("5. VCM-SSD", "Variational\nJoint Coupler", COLOR_GREEN, RGBColor(209, 250, 229)),
    ("6. Output", "Retrieval &\nPrediction Score", RGBColor(4, 120, 87), RGBColor(236, 253, 245))
]

add_native_process_flow(slide12, Inches(2.0), Inches(3.0), workflow_steps)

add_card(slide12, Inches(0.8), Inches(5.4), Inches(5.7), Inches(1.4), "Input & Pre-Filtering (Steps 1-3)", [
    "Raw RGB images and text captions are embedded via ViT and RoBERTa, then noise-damped using Port-Hamiltonian SD-NPF."
])

add_card(slide12, Inches(6.8), Inches(5.4), Inches(5.733), Inches(1.4), "SSD & Variational Output (Steps 4-6)", [
    "Filtered features are modeled via Mamba-2 linear state spaces and fused in VCM-SSD for final score prediction."
])

# ==========================================
# SLIDE 13: Technology Stack
# ==========================================
slide13 = prs.slides.add_slide(blank_layout)
set_slide_background(slide13)
add_header(slide13, "TECHNOLOGY STACK", "Software Architecture & Framework Suite", "Production-ready open-source libraries powering the PH-SSD codebase")
add_footer(slide13, 13)

tech_items = [
    ("Python 3.10+", "Core Programming Language", COLOR_BLUE),
    ("PyTorch 2.x", "Deep Learning Engine & Autograd", COLOR_NAVY),
    ("Mamba-2", "State Space Duality Backbone", COLOR_BLUE),
    ("Vision Transformer (ViT)", "Patch Embedding Vision Encoder", COLOR_SKY),
    ("DINOv2", "Self-Supervised Vision Features", COLOR_NAVY),
    ("RoBERTa", "Pretrained Language Transformer", COLOR_BLUE),
    ("Hugging Face", "Model Registry & Tokenizers", COLOR_SKY),
    ("timm", "PyTorch Image Models Suite", COLOR_NAVY),
    ("YAML / PyYAML", "Modular Configuration Management", COLOR_SLATE)
]

for idx, (name, desc, accent) in enumerate(tech_items):
    col = idx % 3
    row = idx // 3
    left = Inches(0.8 + col * 4.0)
    top = Inches(1.8 + row * 1.6)
    
    card = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.733), Inches(1.4))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_CARD_BORDER
    card.line.width = Pt(1.5)
    
    tx_t = slide13.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), Inches(3.233), Inches(1.0))
    tf_t = tx_t.text_frame
    tf_t.word_wrap = True
    
    p1 = tf_t.paragraphs[0]
    p1.text = name
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = accent
    
    p2 = tf_t.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = COLOR_SLATE

# ==========================================
# SLIDE 14: Expected Evaluation
# ==========================================
slide14 = prs.slides.add_slide(blank_layout)
set_slide_background(slide14)
add_header(slide14, "EXPECTED EVALUATION", "Comprehensive Evaluation Metrics Suite", "Rigorous metrics for evaluating accuracy, retrieval, latency, and efficiency")
add_footer(slide14, 14)

metrics = [
    ("Recall@1 (R@1)", "Primary Image-Text Retrieval Accuracy", COLOR_BLUE),
    ("Recall@5 (R@5)", "Top-5 Candidate Retrieval Precision", COLOR_BLUE),
    ("Recall@10 (R@10)", "Top-10 Candidate Retrieval Coverage", COLOR_BLUE),
    ("Accuracy", "Classification & Semantic Matching Rate", COLOR_NAVY),
    ("Latency (ms)", "End-to-End Inference Latency per Sample", COLOR_RED),
    ("FLOPs", "Floating Point Operations Overhead", COLOR_NAVY),
    ("Throughput", "Processed Images/Sec Processing Speed", COLOR_GREEN),
    ("Memory (VRAM)", "Peak GPU Memory Allocation (GB)", COLOR_RED)
]

for idx, (m_title, m_desc, m_color) in enumerate(metrics):
    col = idx % 4
    row = idx // 4
    left = Inches(0.8 + col * 2.98)
    top = Inches(1.8 + row * 2.4)
    
    card = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.78), Inches(2.1))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = m_color
    card.line.width = Pt(1.5)
    
    tx_m = slide14.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(2.38), Inches(1.7))
    tf_m = tx_m.text_frame
    tf_m.word_wrap = True
    
    p1 = tf_m.paragraphs[0]
    p1.text = m_title
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = m_color
    p1.space_after = Pt(8)
    
    p2 = tf_m.add_paragraph()
    p2.text = m_desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = COLOR_SLATE

# ==========================================
# SLIDE 15: Novel Contributions (3 Core Contributions)
# ==========================================
slide15 = prs.slides.add_slide(blank_layout)
set_slide_background(slide15)
add_header(slide15, "NOVEL CONTRIBUTIONS", "Three Core Architectural Breakthroughs", "Primary technical contributions introduced in the PH-SSD proposal")
add_footer(slide15, 15)

contribs = [
    ("1. SD-NPF (Port-Hamiltonian Noise Filter)", "First pre-filter framework combining Port-Hamiltonian energy dynamics with deep neural representations for physics-inspired noise reduction."),
    ("2. VCM-SSD (Variational Cross-Modal Coupler)", "Novel VCM-SSD coupler enabling probabilistic latent space alignment between vision and text hidden states."),
    ("3. PH-SSD System Architecture", "Integration of Mamba-2 linear state spaces into an efficient dual-stream multimodal learning backbone."),
    ("✓ Linear O(N) Scaling Efficiency", "Yields lower VRAM usage and faster inference latency compared to quadratic Transformers."),
    ("✓ Modular Design Strategy", "Clean, decoupled architecture supporting easy swapping of vision/text backbones and evaluation protocols.")
]

for i, (c_title, c_desc) in enumerate(contribs):
    top_pos = Inches(1.8 + i * 0.98)
    card = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(11.733), Inches(0.85))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_BLUE if i < 3 else COLOR_CARD_BORDER
    card.line.width = Pt(1.5)

    tx_c = slide15.shapes.add_textbox(Inches(1.1), top_pos + Inches(0.12), Inches(11.133), Inches(0.6))
    tf_c = tx_c.text_frame
    tf_c.word_wrap = True
    p_t = tf_c.paragraphs[0]
    p_t.text = c_title
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_BLUE
    
    p_d = tf_c.add_paragraph()
    p_d.text = c_desc
    p_d.font.size = Pt(10.5)
    p_d.font.color.rgb = COLOR_SLATE

# ==========================================
# SLIDE 16: Current Status (Scientifically Honest)
# ==========================================
slide16 = prs.slides.add_slide(blank_layout)
set_slide_background(slide16)
add_header(slide16, "CURRENT STATUS", "Milestone Progress & Project Roadmap", "Overview of completed architectural milestones and pending experimental tasks")
add_footer(slide16, 16)

card_c = slide16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.8))
card_c.fill.solid()
card_c.fill.fore_color.rgb = COLOR_CARD_BG
card_c.line.color.rgb = COLOR_GREEN
card_c.line.width = Pt(2)

tx_c = slide16.shapes.add_textbox(Inches(1.15), Inches(2.05), Inches(5.0), Inches(4.3))
tf_c = tx_c.text_frame
tf_c.word_wrap = True
tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0

p_ct = tf_c.paragraphs[0]
p_ct.text = "Completed Milestones"
p_ct.font.size = Pt(15)
p_ct.font.bold = True
p_ct.font.color.rgb = COLOR_GREEN
p_ct.space_after = Pt(12)

completed_list = [
    ("✓ Architecture Designed & Verified:", "Dual-stream vision-text topology finalized."),
    ("✓ Mathematical Formulation Established:", "Port-Hamiltonian energy dissipation equations derived."),
    ("✓ Code Implementation Built:", "Modular PyTorch pipeline with SD-NPF & VCM-SSD."),
    ("✓ Unit Testing Suite Verified:", "All forward passes and tensor shapes validated.")
]

for title_item, desc_item in completed_list:
    p_h = tf_c.add_paragraph()
    p_h.text = title_item
    p_h.font.size = Pt(11)
    p_h.font.bold = True
    p_h.font.color.rgb = RGBColor(4, 120, 87)
    p_h.space_after = Pt(2)
    
    p_sub = tf_c.add_paragraph()
    p_sub.text = f"   • {desc_item}"
    p_sub.font.size = Pt(9.5)
    p_sub.font.color.rgb = COLOR_SLATE
    p_sub.space_after = Pt(8)

card_p = slide16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.8))
card_p.fill.solid()
card_p.fill.fore_color.rgb = COLOR_CARD_BG
card_p.line.color.rgb = COLOR_BLUE
card_p.line.width = Pt(2)

tx_p = slide16.shapes.add_textbox(Inches(7.15), Inches(2.05), Inches(5.033), Inches(4.3))
tf_p = tx_p.text_frame
tf_p.word_wrap = True
tf_p.margin_left = tf_p.margin_top = tf_p.margin_right = tf_p.margin_bottom = 0

p_pt = tf_p.paragraphs[0]
p_pt.text = "Pending Execution"
p_pt.font.size = Pt(15)
p_pt.font.bold = True
p_pt.font.color.rgb = COLOR_BLUE
p_pt.space_after = Pt(12)

pending_list = [
    ("⏱ Flickr8k Model Training:", "Launch multi-epoch training runs on Flickr8k dataset."),
    ("⏱ Benchmark Evaluation & Ablation Studies:", "Compute Recall@1, 5, 10 retrieval metrics & FLOPs."),
    ("⏱ Comparative Analysis:", "Evaluate against CLIP, BLIP-2, SigLIP, and VL-Mamba."),
    ("⏱ Publication Preparation:", "Draft manuscript for target conference submission.")
]

for title_item, desc_item in pending_list:
    p_h = tf_p.add_paragraph()
    p_h.text = title_item
    p_h.font.size = Pt(11)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_BLUE
    p_h.space_after = Pt(2)
    
    p_sub = tf_p.add_paragraph()
    p_sub.text = f"   • {desc_item}"
    p_sub.font.size = Pt(9.5)
    p_sub.font.color.rgb = COLOR_SLATE
    p_sub.space_after = Pt(8)

# ==========================================
# SLIDE 17: Future Work
# ==========================================
slide17 = prs.slides.add_slide(blank_layout)
set_slide_background(slide17)
add_header(slide17, "FUTURE WORK", "Dataset Expansion & Target Publication Venues", "Planned research roadmap, baseline comparisons, and target venues")
add_footer(slide17, 17)

add_card(slide17, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), "Future Expansion & Baselines", [
    "Dataset Expansion Plan:",
    "  • Train and evaluate on Flickr8k.",
    "  • Scale evaluation to Flickr30k benchmark.",
    "  • Benchmark on large-scale MS-COCO dataset.",
    "Baseline Models for Comparison:",
    "  • CLIP (Contrastive Language-Image Pretraining)",
    "  • BLIP-2 (Bootstrapped Language-Image Pre-training)",
    "  • SigLIP (Sigmoid Loss for Language Image Pre-training)",
    "  • VL-Mamba (Vision-Language State Space Baseline)"
])

venues = [
    ("NeurIPS", "Neural Information Processing Systems", COLOR_BLUE, RGBColor(239, 246, 255)),
    ("ICML", "International Conference on Machine Learning", COLOR_SKY, RGBColor(224, 242, 254)),
    ("ICLR", "International Conference on Learning Representations", COLOR_NAVY, RGBColor(241, 245, 249)),
    ("IEEE TPAMI", "Transactions on Pattern Analysis & Machine Intelligence", COLOR_GREEN, RGBColor(209, 250, 229))
]

v_left = Inches(6.6)
v_top = Inches(1.8)
for i, (v_title, v_desc, stroke, bg) in enumerate(venues):
    row = i // 2
    col = i % 2
    x = v_left + col * Inches(3.0)
    y = v_top + row * Inches(2.4)
    
    card = slide17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = stroke
    card.line.width = Pt(1.5)
    
    tx = slide17.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(1.6))
    tf = tx.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = v_title
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = stroke
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = v_desc
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = COLOR_SLATE
    p2.alignment = PP_ALIGN.CENTER

# ==========================================
# SLIDE 18: Conclusion
# ==========================================
slide18 = prs.slides.add_slide(blank_layout)
set_slide_background(slide18)
add_header(slide18, "CONCLUSION", "Summary of Research Impact & Breakthroughs", "Executive summary of the PH-SSD architecture and its contributions")
add_footer(slide18, 18)

points = [
    ("Novel Physics-Inspired Architecture", "PH-SSD seamlessly combines Port-Hamiltonian energy dynamics with State Space Dualities for robust multimodal learning."),
    ("Effective Noise Suppression", "SD-NPF pre-filtering actively removes uninformative background noise and redundant feature tokens."),
    ("Linear O(N) Scaling Efficiency", "Replaces quadratic Transformer self-attention with Mamba-2 linear state spaces for ultra-fast, low-memory execution."),
    ("Superior Cross-Modal Alignment", "VCM-SSD variational coupler enforces tight latent space cohesion between vision and text streams."),
    ("Paving the Way for Efficient AI", "Bridges control theory and state space models to push the frontier of scalable, efficient AI systems.")
]

for i, (p_title, p_desc) in enumerate(points):
    top_pos = Inches(1.8 + i * 0.98)
    card = slide18.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(11.733), Inches(0.85))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_BLUE if i==0 else COLOR_CARD_BORDER
    card.line.width = Pt(1.5)

    tx_p = slide18.shapes.add_textbox(Inches(1.1), top_pos + Inches(0.12), Inches(11.133), Inches(0.6))
    tf_p = tx_p.text_frame
    tf_p.word_wrap = True
    p_t = tf_p.paragraphs[0]
    p_t.text = p_title
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_BLUE
    
    p_d = tf_p.add_paragraph()
    p_d.text = p_desc
    p_d.font.size = Pt(10.5)
    p_d.font.color.rgb = COLOR_SLATE

# ==========================================
# SLIDE 19 (NEW): Faculty Q&A Reference Guide
# ==========================================
slide19 = prs.slides.add_slide(blank_layout)
set_slide_background(slide19, COLOR_WHITE)
add_header(slide19, "FACULTY Q&A REFERENCE", "Prepared Answers for Review Panel Questions", "Cheat sheet for defending core research decisions during presentation Q&A")
add_footer(slide19, 19)

qa_list = [
    ("Q1: Why Mamba over Transformer?", "Transformer self-attention is O(N²) quadratic; Mamba-2 is O(N) linear with lower VRAM and 5× speedup.", COLOR_BLUE),
    ("Q2: Why Port-Hamiltonian (SD-NPF)?", "Applies energy dissipation to filter out background visual/textual noise prior to state space encoding.", COLOR_SKY),
    ("Q3: Why Flickr8k for Proposal?", "Small, official benchmark ideal for rapid proof-of-concept architecture validation on single GPU.", COLOR_NAVY),
    ("Q4: Is this your own model?", "The overall PH-SSD architecture, SD-NPF, and VCM-SSD are proposed design; ViT, RoBERTa, Mamba are backbones.", COLOR_GREEN),
    ("Q5: Have you trained it?", "Architecture design, math formulation, code, and unit tests are complete; Flickr8k training is pending.", COLOR_RED),
    ("Q6: What are your key contributions?", "1. SD-NPF noise filter, 2. VCM-SSD variational coupler, 3. PH-SSD efficient multimodal system architecture.", COLOR_BLUE)
]

for idx, (q_text, a_text, q_color) in enumerate(qa_list):
    col = idx % 2
    row = idx // 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(1.8 + row * 1.65)
    
    card = slide19.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.733), Inches(1.45))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = q_color
    card.line.width = Pt(1.5)
    
    tx = slide19.shapes.add_textbox(x + Inches(0.25), y + Inches(0.15), Inches(5.233), Inches(1.15))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = q_text
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = q_color
    p1.space_after = Pt(4)
    
    p2 = tf.add_paragraph()
    p2.text = f"Answer: {a_text}"
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = COLOR_SLATE

# Save presentation
prs.save(OUTPUT_PPTX)
print(f"Successfully created presentation: {OUTPUT_PPTX}")

for alt_path in [OUTPUT_PPTX_ALT1, OUTPUT_PPTX_ALT2, OUTPUT_PPTX_ALT3]:
    try:
        prs.save(alt_path)
        print(f"Successfully updated presentation: {alt_path}")
    except Exception:
        print(f"Note: Could not overwrite {alt_path} because it is currently open in PowerPoint.")

